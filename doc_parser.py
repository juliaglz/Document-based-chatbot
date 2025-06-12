import os
import re
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import time
class CleanTextExtractor:
    def __init__(self, file_path):
        self.file_path = file_path
        self.file_extension = os.path.splitext(file_path)[-1].lower()

    @staticmethod
    def clean_text(text):
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r' {2,}', ' ', text)
        return text

    def extract_text(self):
        text_content = []
        if self.file_extension == '.pdf':
            with open(self.file_path, 'rb') as pdf_file:
                reader = PdfReader(pdf_file)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(self.clean_text(page_text))

        elif self.file_extension == '.docx':
            doc = Document(self.file_path)
            text_content = [self.clean_text(para.text) for para in doc.paragraphs if para.text]

        elif self.file_extension == '.pptx':
            prs = Presentation(self.file_path)
            for slide in prs.slides:
                slide_text = [self.clean_text(shape.text) for shape in slide.shapes if hasattr(shape, "text") and shape.text]
                text_content.extend(slide_text)

        elif self.file_extension in ('.xls', '.xlsx'):
            wb = load_workbook(self.file_path)
            text_content = [self.clean_text(str(cell.value)) for sheet in wb for row in sheet.iter_rows() for cell in row if cell.value]

        elif self.file_extension == '.txt':
            with open(self.file_path, 'r', encoding='utf-8') as txt_file:
                text = txt_file.read()
                text_content.append(self.clean_text(text))

        return ' '.join(text_content)

'''
if __name__ == '__main__':
    data_directory = "./wiki2txt"
    start = time.time()
    for filename in os.listdir(data_directory):
        filepath = os.path.join(data_directory, filename)
        extractor = CleanTextExtractor(filepath)
        document_text = extractor.extract_text()
    end = time.time()
    print(end-start)
'''